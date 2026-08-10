"""
Ekstraksi teks dari berkas materi rapat (PDF/DOCX/PPTX/TXT) yang diunggah
saat rapat dibuat, supaya bisa ikut jadi konteks tambahan untuk LLM
(lihat "materi_rapat" di services/summarizer.py) - pola yang sama seperti
catatan_notulis, tapi sumbernya dokumen bukan teks yang diketik manual.

Format lama (.doc/.ppt biner) sengaja TIDAK didukung: butuh dependensi
berat (mis. konversi lewat LibreOffice) untuk manfaat yang kecil. Pesan
errornya eksplisit minta disimpan ulang sebagai .docx/.pptx.
"""
from pathlib import Path

from ..config import settings

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


class UnsupportedMaterialError(Exception):
    pass


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_text(path: Path) -> str:
    """Mengekstrak teks dari berkas materi, dipotong sampai MAX_MATERI_CHARS."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        teks = _extract_pdf(path)
    elif ext == ".docx":
        teks = _extract_docx(path)
    elif ext == ".pptx":
        teks = _extract_pptx(path)
    elif ext == ".txt":
        teks = _extract_txt(path)
    elif ext in (".doc", ".ppt"):
        raise UnsupportedMaterialError(
            f"Format {ext} (versi lama) belum didukung. Simpan ulang sebagai "
            f"{'.docx' if ext == '.doc' else '.pptx'} lalu unggah kembali."
        )
    else:
        raise UnsupportedMaterialError(
            f"Format {ext} tidak didukung untuk materi rapat. "
            f"Format yang didukung: PDF, DOCX, PPTX, TXT."
        )

    teks = teks.strip()
    if len(teks) > settings.MAX_MATERI_CHARS:
        teks = teks[:settings.MAX_MATERI_CHARS] + "\n[...dipotong, materi terlalu panjang...]"
    return teks
